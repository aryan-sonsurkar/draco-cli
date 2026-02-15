from pptx import Presentation
from pptx.util import Inches, Pt
import os
import json
from modules.Ollama.ollama_helper import ask_ollama

ASSETS_DIR = "data/ppt_assets"


def extract_json(text):
    start = text.find("{")
    end = text.rfind("}") + 1
    if start == -1 or end == -1:
        return None
    return text[start:end]


def pick_image(topic):
    topic = topic.lower()

    mapping = {
        "operating system": "os.jpg",
        "os": "os.jpg",
        "network": "networks.jpg",
        "python": "python.jpg",
        "data structure": "datastructures.jpg",
        "database": "database.jpg",
        "ai": "ai.jpg",
    }

    for key in mapping:
        if key in topic:
            path = os.path.join(ASSETS_DIR, mapping[key])
            if os.path.exists(path):
                return path

    return None


def generate_slides(topic):

    prompt = f"""
Create presentation slides in JSON.

Topic: {topic}

Return ONLY valid JSON.

Format:
{{
  "slides": [
    {{
      "title": "Slide title",
      "text": "Slide explanation paragraph"
    }}
  ]
}}

Generate 5–7 slides.
"""

    raw = ask_ollama(prompt)

    print("\n--- AI RAW RESPONSE ---")
    print(raw)
    print("------------------------\n")

    json_block = extract_json(raw)

    if not json_block:
        return []

    try:
        parsed = json.loads(json_block)
        return parsed.get("slides", [])
    except:
        return []


def build_ppt(topic):

    slides_data = generate_slides(topic)

    if not slides_data:
        return "Slide generation failed."

    prs = Presentation()

    image_path = pick_image(topic)

    for slide_info in slides_data:

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.text = slide_info.get("title", "Untitled")

        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(36)
        run.font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.5), Inches(3))
        body_box.text_frame.text = slide_info.get("text", "")

        if image_path:
            slide.shapes.add_picture(
                image_path,
                Inches(2),
                Inches(3.3),
                width=Inches(6)
            )

    filename = topic.replace(" ", "_") + ".pptx"
    prs.save(filename)

    return f"PPT created: {filename}"
